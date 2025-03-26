# ==========================
# Imports da biblioteca padrão
# ==========================
import os
import time
import subprocess
import sys

# ==============================
# 🛠️ Função para verificar e instalar dependências
# ==============================
def instalar_dependencias():
    pacotes = ["requests", "python-pptx", "tqdm", "argostranslate", "keyboard", "pyautogui"]
    for pacote in pacotes:
        try:
            __import__(pacote.replace("-", "_"))  # Ajusta nome para importação
        except ImportError:
            print(f"🔹 Instalando {pacote}...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", pacote])
            time.sleep(2)
    print("✅ Dependências instaladas com sucesso!")

# 1. Verificar e instalar dependências
instalar_dependencias()

# ==========================
# Dependências externas
# ==========================
import requests
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from tqdm import tqdm  # Barra de progresso elegante
import argostranslate.package  # Para gerenciar pacotes de tradução
import keyboard  # Para detectar pressionamento de teclas
import pyautogui  # Para manipular a janela do PowerPoint

# Variável global para a porta do servidor
porta_servidor = 5000
# Variável global para controlar a interrupção
interromper = False
# Idioma de destino predefinido
idioma_destino = "pt"

# ==============================
# 🚀 Funções para gerenciar o servidor LibreTranslate
# ==============================
def verificar_porta_em_uso(porta):
    """Verifica se a porta está em uso."""
    comando = f"netstat -ano | findstr :{porta}"
    resultado = subprocess.run(comando, shell=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return bool(resultado.stdout)

def instalar_modelos_traducao():
    """Garante que o modelo en -> pt esteja instalado."""
    try:
        print("🔹 Verificando e instalando modelo de tradução en -> pt...")
        argostranslate.package.update_package_index()
        available_packages = argostranslate.package.get_available_packages()
        package_to_install = next(
            filter(lambda x: x.from_code == "en" and x.to_code == "pt", available_packages), None
        )
        if package_to_install:
            print(f"🔹 Baixando e instalando modelo {package_to_install}...")
            argostranslate.package.install_from_path(package_to_install.download())
            print("✅ Modelo de tradução en -> pt instalado com sucesso!")
        else:
            print("❌ Modelo de tradução en -> pt não encontrado no índice!")
            return False
        return True
    except Exception as e:
        print(f"❌ Erro ao instalar modelo de tradução: {e}")
        return False

def instalar_dependencias_libretranslate():
    """Instala dependências específicas do LibreTranslate."""
    try:
        print("🔹 Instalando dependências do LibreTranslate...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "Flask", "flask_cors", "requests", "argostranslate"])
        print("✅ Dependências do LibreTranslate instaladas com sucesso!")
        return True
    except subprocess.CalledProcessError as e:
        print(f"❌ Erro ao instalar dependências do LibreTranslate: {e}")
        return False

def instalar_libretranslate():
    """Instala o servidor LibreTranslate baixando do repositório GitHub e configurando as dependências."""
    try:
        print("🔹 Instalando o LibreTranslate...")
        if not os.path.exists("LibreTranslate"):
            subprocess.run(["git", "clone", "https://github.com/LibreTranslate/LibreTranslate.git"], check=True)
        os.chdir("LibreTranslate")

        # Instalar as dependências do LibreTranslate
        if instalar_dependencias_libretranslate():
            print("✅ LibreTranslate instalado com sucesso!")
        else:
            print("❌ Falha ao instalar as dependências do LibreTranslate!")
            return False
        return True
    except subprocess.CalledProcessError as e:
        print(f"❌ Erro ao instalar o LibreTranslate: {e}")
        return False

def iniciar_servidor():
    """Inicia o servidor LibreTranslate e verifica se está funcional."""
    global porta_servidor, idioma_destino
    if verificar_porta_em_uso(porta_servidor):
        print(f"⚠️ A porta {porta_servidor} já está em uso! Tentando com a porta 5001...")
        porta_servidor = 5001

    try:
        # Baixar e instalar o LibreTranslate se necessário
        if not instalar_libretranslate():
            return None
        
        # Iniciar o servidor
        print(f"🔄 Iniciando o servidor LibreTranslate na porta {porta_servidor}...")
        servidor = subprocess.Popen(
            ["python", "main.py", "--port", str(porta_servidor), "--load-only", "en,pt"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        time.sleep(20)  # Tempo para carregar os modelos

        # Verificar se o servidor está respondendo
        response = requests.get(f"http://localhost:{porta_servidor}", timeout=5)
        if response.status_code == 200:
            print("✅ Servidor iniciado com sucesso!")
            # Verificar idiomas suportados
            langs_otimizado = requests.get(f"http://localhost:{porta_servidor}/languages", timeout=5)
            langs = langs_otimizado.json()
            # Verificar se o idioma predefinido 'pt' está disponível como alvo
            pt_supported = any("pt" in lang["targets"] for lang in langs)
            if not pt_supported:
                print("❌ Português (pt) não está suportado pelo servidor!")
                servidor.terminate()
                return None
            
            # Exibir apenas o idioma predefinido e solicitar confirmação
            print(f"🔹 Idioma de destino predefinido: Português (pt)")
            print("🔹 Deseja prosseguir com este idioma? (s/n) ou alterar (a): ", end="")
            escolha = input().strip().lower()
            print(f"Resposta do usuário: {escolha}")

            if escolha == "a":
                # Mostrar todos os idiomas disponíveis para escolha
                print("🔹 Idiomas suportados pelo servidor:")
                for lang in langs:
                    print(f"  - {lang['code']}: {lang['name']} (Targets: {', '.join(lang['targets'])})")
                print("🔹 Digite o código do idioma desejado (ex.: 'es' para espanhol): ", end="")
                novo_idioma = input().strip().lower()
                print(f"Novo idioma escolhido: {novo_idioma}")
                # Verificar se o novo idioma é suportado
                idioma_valido = any(novo_idioma in lang["targets"] for lang in langs)
                if idioma_valido:
                    idioma_destino = novo_idioma
                    print(f"✅ Idioma alterado para: {novo_idioma}")
                else:
                    print(f"❌ Idioma '{novo_idioma}' não suportado! Mantendo 'pt' como padrão.")
            elif escolha != "s":
                print("❌ Operação cancelada pelo usuário.")
                servidor.terminate()
                return None

            return servidor
        else:
            print(f"❌ Servidor não respondeu na porta {porta_servidor}. Status: {response.status_code}")
            servidor.terminate()
            return None

    except subprocess.CalledProcessError as e:
        print(f"❌ Erro ao iniciar o servidor: {e}")
        return None
    except requests.exceptions.RequestException as e:
        print(f"❌ Servidor não está acessível: {e}")
        erro_servidor = servidor.stderr.read()
        if erro_servidor:
            print(f"⚠️ Saída de erro do servidor: {erro_servidor}")
        servidor.terminate()
        return None

# ==============================
# 🌍 Função para traduzir textos usando LibreTranslate
# ==============================
def traduzir_texto(texto, de="en", para=idioma_destino):
    """Traduz texto do inglês para o idioma de destino usando o servidor LibreTranslate."""
    global interromper
    if interromper:
        return texto  # Retorna o texto original se a operação foi interrompida
    url = f"http://localhost:{porta_servidor}/translate"
    payload = {"q": texto, "source": de, "target": para, "format": "text"}
    
    try:
        response = requests.post(url, json=payload, timeout=10)
        if response.status_code == 200:
            return response.json().get("translatedText", texto)
        else:
            print(f"⚠️ Erro na tradução: {response.text}")
            return texto
    except requests.exceptions.RequestException as e:
        print(f"❌ Erro ao conectar ao LibreTranslate: {e}")
        return texto

# ==============================
# 📂 Abrir janela para selecionar arquivo PowerPoint
# ==============================
def selecionar_arquivo():
    """Abre uma janela para selecionar um arquivo PowerPoint."""
    root = tk.Tk()
    root.withdraw()  # Esconde a janela raiz
    root.attributes('-topmost', True)  # Força a janela a ficar no topo
    caminho_arquivo = filedialog.askopenfilename(
        title="Selecione um arquivo PowerPoint",
        filetypes=[("PowerPoint", "*.pptx")]
    )
    root.attributes('-topmost', False)  # Remove o estado "topmost" após a seleção
    root.destroy()  # Fecha a instância do Tkinter
    return caminho_arquivo

# ==============================
# 🔄 Traduzir o conteúdo do PowerPoint mantendo a formatação
# ==============================
def traduzir_powerpoint(caminho_ppt, servidor):
    """Traduz o conteúdo de um arquivo PowerPoint do inglês para o idioma de destino, preservando formatação."""
    global interromper
    if not caminho_ppt:
        print("❌ Nenhum arquivo selecionado.")
        return
    
    print(f"📂 Arquivo carregado: {caminho_ppt}")
    print("ℹ️ Pressione 'q' a qualquer momento para abortar a tradução.")
    prs = Presentation(caminho_ppt)
    
    # Contar o total de runs para a barra de progresso
    total_runs = sum(
        len(paragraph.runs)
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text_frame") and shape.text_frame.text.strip()
        for paragraph in shape.text_frame.paragraphs
        if paragraph.runs
    )
    traduzidos = 0

    with tqdm(total=total_runs, desc="Traduzindo", unit="segmento") as pbar:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if interromper:
                                print("⚠️ Tradução interrompida pelo usuário!")
                                servidor.terminate()
                                print("🛑 Servidor encerrado.")
                                return
                            if run.text.strip():  # Traduzir apenas runs com texto não vazio
                                run.text = traduzir_texto(run.text, de="en", para=idioma_destino)
                                traduzidos += 1
                                pbar.update(1)

    if not interromper:
        novo_nome = caminho_ppt.replace(".pptx", f"_traduzido_{idioma_destino}.pptx")
        prs.save(novo_nome)
        print(f"✅ Tradução concluída! Arquivo salvo como: {novo_nome}")

# ==============================
# 🚀 Função para trazer a janela do PowerPoint para o primeiro plano
# ==============================
def trazer_ppt_para_primeiro_plano():
    """Força a janela do PowerPoint para o primeiro plano."""
    try:
        pyautogui.getWindowsWithTitle("PowerPoint")[0].activate()  # Ativa a janela com "PowerPoint" no título
    except IndexError:
        print("❌ Não foi possível encontrar uma janela do PowerPoint.")

# ==============================
# 🚀 Fluxo principal do programa
# ==============================
if __name__ == "__main__":
    # Configurar o listener para a tecla 'q'
    def on_q_press(event):
        global interromper
        interromper = True
        print("\n⚠️ Tecla 'q' pressionada. Abortando operação...")

    keyboard.on_press_key("q", on_q_press)

    # Iniciar o servidor LibreTranslate
    print("🔄 Iniciando servidor LibreTranslate...")
    servidor = iniciar_servidor()
    
    if servidor is None:
        print("❌ Não foi possível iniciar o servidor. Encerrando o programa.")
        sys.exit(1)

    # Perguntar ao usuário se deseja traduzir
    print("🔹 Deseja iniciar a tradução? (s/n): ", end="")
    iniciar = input().strip().lower()
    print(f"Resposta do usuário: {iniciar}")

    if iniciar == "s":
        print("📂 Abrindo janela para selecionar PowerPoint...")
        caminho_ppt = selecionar_arquivo()
        print(f"Arquivo selecionado: {caminho_ppt}")
        
        # Tente trazer a janela do PowerPoint para o primeiro plano
        trazer_ppt_para_primeiro_plano()
        
        traduzir_powerpoint(caminho_ppt, servidor)
    else:
        print("❌ Tradução cancelada.")

    # Encerrar o servidor
    if not interromper:
        print("🛑 Encerrando servidor LibreTranslate...")
        servidor.terminate()
    print("✅ Programa finalizado.")
